from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import json
import base64
import re
import io
import zipfile
import tempfile
import shutil
import os
from color_handler import ColorHandler
from text_handler import TextHandler
from shape_handler import ShapeHandler

class PPTXFabricConverter:
    def __init__(self):
        self.debug = True  # Enable debug logging
        
    def _fix_invalid_namespace_uri(self, xml_content):
        """Fix invalid namespace URIs in XML content"""
        # Fix backslashes in namespace URIs
        xml_content = re.sub(r'xmlns:[^=]+="[^"]+\\"', lambda m: m.group(0)[:-1] + '"', xml_content)
        
        # Fix specific problematic namespaces
        namespace_fixes = {
            'http://schemas.microsoft.com/office/drawing/2014/main\\': 'http://schemas.microsoft.com/office/drawing/2014/main',
            'http://schemas.microsoft.com/office/powerpoint/2010/main\\': 'http://schemas.microsoft.com/office/powerpoint/2010/main',
            'http://schemas.microsoft.com/office/word/2010/wordprocessingDrawing\\': 'http://schemas.microsoft.com/office/word/2010/wordprocessingDrawing'
        }
        
        for invalid, valid in namespace_fixes.items():
            xml_content = xml_content.replace(invalid, valid)
            
        return xml_content
        
    def _preprocess_pptx_file(self, pptx_file):
        """Preprocess PPTX file to fix any XML issues"""
        # Create a temporary file to store the processed PPTX
        temp_dir = tempfile.mkdtemp()
        try:
            # Save the uploaded file to a temporary location
            temp_pptx_path = os.path.join(temp_dir, 'original.pptx')
            pptx_file.save(temp_pptx_path)
            
            # Extract the PPTX file
            with zipfile.ZipFile(temp_pptx_path, 'r') as zip_ref:
                zip_ref.extractall(temp_dir)
            
            # Process XML files
            for root, _, files in os.walk(temp_dir):
                for file in files:
                    if file.endswith('.xml'):
                        file_path = os.path.join(root, file)
                        try:
                            with open(file_path, 'r', encoding='utf-8') as f:
                                content = f.read()
                            
                            # Fix invalid namespace URIs
                            fixed_content = self._fix_invalid_namespace_uri(content)
                            
                            with open(file_path, 'w', encoding='utf-8') as f:
                                f.write(fixed_content)
                        except Exception as e:
                            print(f"Error processing file {file_path}: {e}")
                            continue
            
            # Create a new PPTX file
            new_pptx_path = os.path.join(temp_dir, 'fixed.pptx')
            with zipfile.ZipFile(new_pptx_path, 'w', zipfile.ZIP_DEFLATED) as zip_ref:
                for root, _, files in os.walk(temp_dir):
                    for file in files:
                        if file not in ['original.pptx', 'fixed.pptx']:
                            file_path = os.path.join(root, file)
                            arcname = os.path.relpath(file_path, temp_dir)
                            zip_ref.write(file_path, arcname)
            
            return new_pptx_path
            
        except Exception as e:
            print(f"Error preprocessing PPTX file: {e}")
            return temp_pptx_path  # Return original file if preprocessing fails
        
        finally:
            try:
                # Clean up the original file
                if os.path.exists(temp_pptx_path):
                    os.remove(temp_pptx_path)
            except Exception as e:
                print(f"Error cleaning up temporary file: {e}")
        
    def pptx_to_fabric(self, pptx_file):
        """Convert PowerPoint file to Fabric.js JSON format"""
        processed_pptx = None
        try:
            # Preprocess the PPTX file
            processed_pptx = self._preprocess_pptx_file(pptx_file)
            
            # Load presentation
            prs = Presentation(processed_pptx)
            
            # Initialize handlers
            color_handler = ColorHandler(prs)
            text_handler = TextHandler(color_handler)
            shape_handler = ShapeHandler(color_handler, text_handler)
            
            # Process each slide
            slides_data = []
            for slide_index, slide in enumerate(prs.slides):
                if self.debug:
                    print(f"\nProcessing slide {slide_index + 1}")
                
                slide_data = {
                    "objects": [],
                    "width": prs.slide_width / 12700,  # Convert EMU to points
                    "height": prs.slide_height / 12700,
                    "slideNumber": slide_index + 1,
                    "background": self._get_slide_background(slide, color_handler)
                }
                
                # Process shapes
                for shape in slide.shapes:
                    if self.debug:
                        print(f"\nShape type: {shape.shape_type}")
                        print(f"Shape name: {shape.name}")
                        if hasattr(shape, 'text'):
                            print(f"Shape text: {shape.text}")
                    
                    shape_data = shape_handler.process_shape(shape)
                    if shape_data:
                        if shape_data["type"] == "group":
                            slide_data["objects"].extend(shape_data["objects"])
                        else:
                            slide_data["objects"].append(shape_data)
                
                slides_data.append(slide_data)
            
            return slides_data
            
        except Exception as e:
            print(f"Error converting PPTX to Fabric: {e}")
            raise
        finally:
            # Clean up the processed file
            if processed_pptx and os.path.exists(processed_pptx):
                try:
                    os.remove(processed_pptx)
                    # Clean up the parent directory
                    temp_dir = os.path.dirname(processed_pptx)
                    if os.path.exists(temp_dir):
                        shutil.rmtree(temp_dir, ignore_errors=True)
                except Exception as e:
                    print(f"Error cleaning up temporary file: {e}")
    
    def fabric_to_pptx(self, fabric_data, template_pptx=None):
        """Convert Fabric.js JSON format back to PowerPoint"""
        try:
            # Create new presentation or use template
            prs = Presentation(template_pptx) if template_pptx else Presentation()
            
            # Set slide size if not using template
            if not template_pptx:
                prs.slide_width = int(fabric_data[0]["width"] * 12700)
                prs.slide_height = int(fabric_data[0]["height"] * 12700)
            
            # Process each slide
            for slide_data in fabric_data:
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
                
                # Set background
                if "background" in slide_data:
                    self._set_slide_background(slide, slide_data["background"])
                
                # Process objects
                for obj in slide_data["objects"]:
                    self._create_shape_from_fabric(slide, obj)
            
            return prs
            
        except Exception as e:
            print(f"Error converting Fabric to PPTX: {e}")
            raise
    
    def _get_slide_background(self, slide, color_handler):
        """Extract slide background properties"""
        try:
            if hasattr(slide, 'background') and slide.background.fill:
                fill_color = color_handler.get_shape_color(slide.background)
                if fill_color:
                    return {
                        "type": "rect",
                        "left": 0,
                        "top": 0,
                        "width": "100%",
                        "height": "100%",
                        "fill": fill_color,
                        "selectable": False
                    }
            return None
        except Exception as e:
            print(f"Error getting slide background: {e}")
            return None
    
    def _set_slide_background(self, slide, background):
        """Set slide background from Fabric.js data"""
        try:
            if background and "fill" in background:
                fill = slide.background.fill
                color = self._parse_color(background["fill"])
                if color:
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(*color)
        except Exception as e:
            print(f"Error setting slide background: {e}")
    
    def _create_shape_from_fabric(self, slide, obj):
        """Create PowerPoint shape from Fabric.js object"""
        try:
            shape = None
            
            if obj["type"] == "textbox":
                shape = slide.shapes.add_textbox(
                    Inches(obj["left"] / 72),
                    Inches(obj["top"] / 72),
                    Inches(obj["width"] / 72),
                    Inches(obj["height"] / 72)
                )
                self._set_text_properties(shape.text_frame, obj)
                
            elif obj["type"] == "rect":
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(obj["left"] / 72),
                    Inches(obj["top"] / 72),
                    Inches(obj["width"] / 72),
                    Inches(obj["height"] / 72)
                )
                self._set_shape_properties(shape, obj)
                
            elif obj["type"] == "path":
                # Convert path to freeform shape
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.FREEFORM,
                    Inches(obj["left"] / 72),
                    Inches(obj["top"] / 72),
                    Inches(obj["width"] / 72),
                    Inches(obj["height"] / 72)
                )
                self._set_shape_properties(shape, obj)
                
            elif obj["type"] == "image":
                # Convert base64 image to bytes
                if "src" in obj and obj["src"].startswith("data:image"):
                    img_data = base64.b64decode(obj["src"].split(",")[1])
                    shape = slide.shapes.add_picture(
                        io.BytesIO(img_data),
                        Inches(obj["left"] / 72),
                        Inches(obj["top"] / 72),
                        Inches(obj["width"] / 72),
                        Inches(obj["height"] / 72)
                    )
            
            if shape and "angle" in obj:
                shape.rotation = float(obj["angle"])
            
        except Exception as e:
            print(f"Error creating shape from Fabric object: {e}")
    
    def _set_text_properties(self, text_frame, obj):
        """Set text properties from Fabric.js object"""
        try:
            if "text" not in obj:
                return
                
            paragraph = text_frame.paragraphs[0]
            run = paragraph.add_run()
            run.text = obj["text"]
            
            font = run.font
            if "fontSize" in obj:
                font.size = Pt(obj["fontSize"])
            if "fontFamily" in obj:
                font.name = obj["fontFamily"]
            if "fill" in obj:
                color = self._parse_color(obj["fill"])
                if color:
                    font.color.rgb = RGBColor(*color)
            if "fontWeight" in obj:
                font.bold = obj["fontWeight"] == "bold"
            if "fontStyle" in obj:
                font.italic = obj["fontStyle"] == "italic"
            if "textAlign" in obj:
                paragraph.alignment = self._get_alignment(obj["textAlign"])
                
        except Exception as e:
            print(f"Error setting text properties: {e}")
    
    def _set_shape_properties(self, shape, obj):
        """Set shape properties from Fabric.js object"""
        try:
            if "fill" in obj:
                fill = shape.fill
                color = self._parse_color(obj["fill"])
                if color:
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(*color)
                    
            if "stroke" in obj:
                line = shape.line
                color = self._parse_color(obj["stroke"])
                if color:
                    line.color.rgb = RGBColor(*color)
                if "strokeWidth" in obj:
                    line.width = Pt(obj["strokeWidth"])
                    
        except Exception as e:
            print(f"Error setting shape properties: {e}")
    
    def _parse_color(self, color):
        """Parse color string to RGB tuple"""
        try:
            if isinstance(color, str):
                if color.startswith('#'):
                    # Hex color
                    color = color.lstrip('#')
                    return tuple(int(color[i:i+2], 16) for i in (0, 2, 4))
                elif color.startswith('rgb'):
                    # RGB color
                    return tuple(map(int, re.findall(r'\d+', color)))
            return None
        except Exception as e:
            print(f"Error parsing color: {e}")
            return None
    
    def _get_alignment(self, align):
        """Convert Fabric.js alignment to PowerPoint alignment"""
        from pptx.enum.text import PP_ALIGN
        align_map = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT,
            'justify': PP_ALIGN.JUSTIFY
        }
        return align_map.get(align.lower(), PP_ALIGN.LEFT) 