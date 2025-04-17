import os
import uuid
from datetime import datetime
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL, MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
import math
import base64
import io
import zipfile
import tempfile
import shutil
import re
from color_handler import ColorHandler
from text_handler import TextHandler
from shape_handler import ShapeHandler

def handle_picture(shape):
    """Handle picture shapes"""
    try:
        image_data = None
        if hasattr(shape, 'image') and shape.image:
            image_blob = shape.image.blob
            image_type = shape.image.content_type.split('/')[-1]  # e.g., 'jpeg', 'png'
            image_data = base64.b64encode(image_blob).decode()
        
        opacity = 1
        if hasattr(shape, 'fill') and shape.fill:
            opacity = 1 - (shape.fill.transparency or 0)
        
        return {
            "type": "image",
            "src": f"data:image/{image_type};base64,{image_data}" if image_data else None,
            "opacity": opacity
        }
    except Exception as e:
        print(f"Error processing image: {e}")
        return None

def get_color_value(color):
    """Extract color value safely from a shape color"""
    try:
        if hasattr(color, 'rgb') and color.rgb:
            # Direct RGB color
            return f'#{color.rgb[0]:02x}{color.rgb[1]:02x}{color.rgb[2]:02x}'
        elif hasattr(color, 'theme_color'):
            # Get color from theme if available
            if hasattr(color._theme, 'theme_elements') and color._theme.theme_elements:
                theme = color._theme.theme_elements
                if theme.clrScheme:
                    scheme = theme.clrScheme
                    if str(color.theme_color) == 'ACCENT_1':
                        return get_scheme_color(scheme.accent1)
                    elif str(color.theme_color) == 'ACCENT_2':
                        return get_scheme_color(scheme.accent2)
                    elif str(color.theme_color) == 'ACCENT_3':
                        return get_scheme_color(scheme.accent3)
                    elif str(color.theme_color) == 'ACCENT_4':
                        return get_scheme_color(scheme.accent4)
                    elif str(color.theme_color) == 'ACCENT_5':
                        return get_scheme_color(scheme.accent5)
                    elif str(color.theme_color) == 'ACCENT_6':
                        return get_scheme_color(scheme.accent6)
                    elif str(color.theme_color) == 'BACKGROUND_1':
                        return get_scheme_color(scheme.bg1)
                    elif str(color.theme_color) == 'BACKGROUND_2':
                        return get_scheme_color(scheme.bg2)
                    elif str(color.theme_color) == 'TEXT_1':
                        return get_scheme_color(scheme.tx1)
                    elif str(color.theme_color) == 'TEXT_2':
                        return get_scheme_color(scheme.tx2)
            
            # If theme color not found, try to get RGB
            if hasattr(color, 'rgb') and color.rgb:
                return f'#{color.rgb[0]:02x}{color.rgb[1]:02x}{color.rgb[2]:02x}'
    except Exception as e:
        print(f"Error getting color value: {e}")
    return None  # Return None instead of default color

def get_scheme_color(scheme_color):
    """Extract color from scheme color element"""
    try:
        if hasattr(scheme_color, 'srgbClr'):
            # Direct sRGB color
            return f'#{scheme_color.srgbClr.val}'
        elif hasattr(scheme_color, 'sysClr'):
            # System color
            return f'#{scheme_color.sysClr.lastClr}'
    except Exception as e:
        print(f"Error getting scheme color: {e}")
    return None

def get_shape_fill_info(shape):
    """Get comprehensive fill information for shapes"""
    if not hasattr(shape, 'fill'):
        return None

    try:
        fill = shape.fill
        if fill is None or fill.type is None:
            return None
            
        if fill.type == MSO_FILL.SOLID:
            if hasattr(fill, 'fore_color') and fill.fore_color:
                color = get_color_value(fill.fore_color)
                if color:
                    return {
                        'type': 'solid',
                        'value': color
                    }
        elif fill.type == MSO_FILL.GRADIENT:
            # Extract gradient information from XML
            if hasattr(shape, 'element'):
                grad_fill = shape.element.find('.//a:gradFill', 
                    {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                if grad_fill is not None:
                    return extract_gradient_info(grad_fill)
    except Exception as e:
        print(f"Error getting fill info: {e}")
    
    return None

def extract_gradient_info(grad_fill):
    """Extract gradient information from XML element"""
    try:
        gradient_info = {'type': 'gradient', 'value': {'type': 'linear', 'colorStops': {}}}
        
        # Get gradient type
        path = grad_fill.find('.//a:path', 
            {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
        if path is not None:
            path_type = path.get('path', 'linear')
            gradient_info['value']['type'] = path_type
        
        # Get gradient stops
        gs_list = grad_fill.findall('.//a:gs', 
            {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
        
        for gs in gs_list:
            pos = float(gs.get('pos', '0')) / 100000
            
            # Get color from different possible sources
            color = None
            for color_type in ['srgbClr', 'schemeClr', 'sysClr']:
                color_elem = gs.find(f'.//a:{color_type}', 
                    {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                if color_elem is not None:
                    if color_type == 'srgbClr':
                        color = f'#{color_elem.get("val")}'
                        break
                    elif color_type == 'schemeClr':
                        # Get color from theme
                        theme_color = color_elem.get('val')
                        # You would need to implement getting the actual color from theme
                        break
                    elif color_type == 'sysClr':
                        color = f'#{color_elem.get("lastClr")}'
                        break
            
            if color:
                gradient_info['value']['colorStops'][str(pos)] = color
        
        return gradient_info if gradient_info['value']['colorStops'] else None
        
    except Exception as e:
        print(f"Error extracting gradient info: {e}")
        return None

def get_line_properties(shape):
    """Get line properties of a shape"""
    try:
        if hasattr(shape, 'line'):
            line = shape.line
            props = {}
            
            if line.color:
                color = get_color_value(line.color)
                if color:
                    props['stroke'] = color
            
            if hasattr(line, 'width'):
                props['strokeWidth'] = line.width / 12700  # Convert to points
                
            return props
    except Exception as e:
        print(f"Error getting line properties: {e}")
    
    return {}

def get_text_properties(shape):
    text_props = {
        'paragraphs': []
    }
    
    if not hasattr(shape, 'text_frame'):
        return text_props
    
    # Process each paragraph in the shape
    for paragraph in shape.text_frame.paragraphs:
        para_props = {
            'text': paragraph.text or '',
            'align': 'left'  # Default alignment
        }
        
        # Font properties from the first run (assuming consistent formatting)
        if paragraph.runs:
            font = paragraph.runs[0].font
            if hasattr(font, 'size') and font.size:
                para_props['fontSize'] = font.size.pt
            if hasattr(font, 'name') and font.name:
                para_props['fontFamily'] = font.name
            if hasattr(font, 'bold'):
                para_props['fontWeight'] = 'bold' if font.bold else 'normal'
            if hasattr(font, 'italic'):
                para_props['fontStyle'] = 'italic' if font.italic else 'normal'
            if hasattr(font, 'color') and font.color and hasattr(font.color, 'rgb'):
                rgb = font.color.rgb
                para_props['fill'] = f'#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}'
        
        text_props['paragraphs'].append(para_props)
    
    return text_props

def get_shape_path(shape):
    """Extract shape path data for complex shapes"""
    try:
        # For FREEFORM shapes, get the vertices
        if shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
            path = []
            vertices = shape.vertices
            for i, (x, y) in enumerate(vertices):
                if i == 0:
                    path.append(f'M {x} {y}')
                else:
                    path.append(f'L {x} {y}')
            path.append('Z')  # Close the path
            return [' '.join(path)]
            
        # Original path extraction for other shapes
        if hasattr(shape, 'element'):
            # Get shape geometry
            spTree = shape.element
            path_list = []
            
            # Find all path elements
            for sp in spTree.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}sp'):
                for spPr in sp.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}spPr'):
                    for custGeom in spPr.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}custGeom'):
                        for pathLst in custGeom.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}pathLst'):
                            for path in pathLst.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}path'):
                                current_path = []
                                
                                # Process move commands
                                for moveTo in path.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}moveTo'):
                                    for pt in moveTo.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}pt'):
                                        x = float(pt.get('x'))
                                        y = float(pt.get('y'))
                                        current_path.append(f'M {x} {y}')
                                
                                # Process line commands
                                for lineTo in path.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}lnTo'):
                                    for pt in lineTo.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}pt'):
                                        x = float(pt.get('x'))
                                        y = float(pt.get('y'))
                                        current_path.append(f'L {x} {y}')
                                
                                # Process cubic bezier curves
                                for cubicBezTo in path.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}cubicBezTo'):
                                    points = []
                                    for pt in cubicBezTo.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}pt'):
                                        x = float(pt.get('x'))
                                        y = float(pt.get('y'))
                                        points.append((x, y))
                                    if len(points) == 3:
                                        current_path.append(f'C {points[0][0]} {points[0][1]} {points[1][0]} {points[1][1]} {points[2][0]} {points[2][1]}')
                                
                                # Process arc commands
                                for arcTo in path.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}arcTo'):
                                    # Convert arc parameters to bezier curves
                                    # This is a simplified version - you might need more complex arc handling
                                    rx = float(arcTo.get('rx', 0))
                                    ry = float(arcTo.get('ry', 0))
                                    angle = float(arcTo.get('angle', 0))
                                    current_path.append(f'A {rx} {ry} {angle} 0 1')
                                
                                # Close path if needed
                                if path.get('w') == '1':
                                    current_path.append('Z')
                                
                                path_list.append(' '.join(current_path))
                                
            return path_list if path_list else None

    except Exception as e:
        print(f"Error extracting shape path: {e}")
        return None

def get_freeform_path(shape):
    """Extract path data from a freeform shape"""
    try:
        path_elem = shape.element.find('.//a:path', 
            {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
        if path_elem is None:
            return None

        path_commands = []
        current_path = []
        
        for child in path_elem:
            tag = child.tag.split('}')[-1]
            if tag == 'moveTo':
                if current_path:
                    path_commands.append(current_path)
                current_path = ['M']
                pt = child.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}pt')
                x = float(pt.get('x')) / 12700
                y = float(pt.get('y')) / 12700
                current_path.extend([x, y])
            elif tag == 'lnTo':
                pt = child.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}pt')
                x = float(pt.get('x')) / 12700
                y = float(pt.get('y')) / 12700
                current_path.extend(['L', x, y])
            elif tag == 'cubicBezTo':
                pts = child.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}pt')
                if len(pts) == 3:
                    current_path.append('C')
                    for pt in pts:
                        x = float(pt.get('x')) / 12700
                        y = float(pt.get('y')) / 12700
                        current_path.extend([x, y])
            elif tag == 'close':
                current_path.append('Z')
                path_commands.append(current_path)
                current_path = []

        if current_path:
            path_commands.append(current_path)

        # Convert path commands to SVG path string
        svg_paths = []
        for commands in path_commands:
            path = ''
            i = 0
            while i < len(commands):
                if isinstance(commands[i], str):
                    path += commands[i] + ' '
                    i += 1
                else:
                    path += f"{commands[i]},{commands[i+1]} "
                    i += 2
            svg_paths.append(path.strip())

        return svg_paths

    except Exception as e:
        print(f"Error getting path data: {e}")
        return None

def process_group_shape(group_shape):
    """Process all shapes within a group"""
    shapes_data = []
    
    for shape in group_shape.shapes:
        shape_data = process_shape(shape)
        if shape_data:
            shapes_data.append(shape_data)
    
    return shapes_data

def process_shape(shape):
    """Process individual shape with all its properties"""
    try:
        base_data = {
            "left": shape.left / 12700,
            "top": shape.top / 12700,
            "width": shape.width / 12700,
            "height": shape.height / 12700,
            "angle": shape.rotation if hasattr(shape, 'rotation') else 0
        }

        # Get shape fill and line properties
        fill_info = get_shape_fill_info(shape)
        line_props = get_line_properties(shape)

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            group_shapes = []
            for child in shape.shapes:
                child_data = process_shape(child)
                if child_data:
                    group_shapes.append(child_data)
            return {
                **base_data,
                "type": "group",
                "objects": group_shapes
            }
        
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pic_data = handle_picture(shape)
            if pic_data:
                return {**base_data, **pic_data}
            return None

        elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            text_data = {
                "type": "textbox",
                "text": shape.text,
                "fontFamily": "Arial",
                "fontSize": 12,
                "textAlign": "center"
            }
            
            # Get text properties
            if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
                para = shape.text_frame.paragraphs[0]
                if para.runs:
                    run = para.runs[0]
                    if hasattr(run.font, 'size') and run.font.size:
                        text_data["fontSize"] = run.font.size.pt
                    if hasattr(run.font, 'name') and run.font.name:
                        text_data["fontFamily"] = run.font.name
                    if hasattr(run.font, 'color') and run.font.color:
                        color = get_color_value(run.font.color)
                        if color:
                            text_data["fill"] = color
                    if hasattr(para, 'alignment'):
                        text_data["textAlign"] = str(para.alignment).lower()
            
            # Add fill and line properties
            if fill_info:
                text_data["backgroundColor"] = fill_info['value'] if fill_info['type'] == 'solid' else None
            
            text_data.update(line_props)
            
            return {**base_data, **text_data}

        elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            shape_data = {
                **base_data,
                "type": "shape",
                **line_props
            }
            
            if fill_info:
                if fill_info['type'] == 'solid':
                    shape_data["fill"] = fill_info['value']
                else:
                    shape_data["gradient"] = fill_info['value']
            
            return shape_data

        elif shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
            path_data = get_freeform_path(shape)
            if path_data:
                shape_data = {
                    **base_data,
                    "type": "path",
                    "path": path_data,
                    **line_props
                }
                
                if fill_info:
                    if fill_info['type'] == 'solid':
                        shape_data["fill"] = fill_info['value']
                    else:
                        shape_data["gradient"] = fill_info['value']
                
                return shape_data

        return None

    except Exception as e:
        print(f"Error processing shape: {e}")
        return None

def fix_invalid_namespace_uri(xml_content):
    """Fix invalid namespace URIs in XML content"""
    xml_content = re.sub(r'xmlns:[^=]+="[^"]+\\"', lambda m: m.group(0)[:-1] + '"', xml_content)
    xml_content = xml_content.replace('http://schemas.microsoft.com/office/drawing/2014/main\\', 
                                    'http://schemas.microsoft.com/office/drawing/2014/main')
    return xml_content

def preprocess_pptx_file(pptx_file):
    """Preprocess PPTX file to fix any XML issues"""
    temp_dir = tempfile.mkdtemp()
    try:
        # Extract the PPTX file
        with zipfile.ZipFile(pptx_file, 'r') as zip_ref:
            zip_ref.extractall(temp_dir)
        
        # Process XML files
        for root, _, files in os.walk(temp_dir):
            for file in files:
                if file.endswith('.xml'):
                    file_path = os.path.join(root, file)
                    with open(file_path, 'r', encoding='utf-8') as f:
                        content = f.read()
                    
                    # Fix invalid namespace URIs
                    fixed_content = fix_invalid_namespace_uri(content)
                    
                    with open(file_path, 'w', encoding='utf-8') as f:
                        f.write(fixed_content)
        
        # Create a new PPTX file
        new_pptx_path = os.path.join(temp_dir, 'fixed.pptx')
        with zipfile.ZipFile(new_pptx_path, 'w', zipfile.ZIP_DEFLATED) as zip_ref:
            for root, _, files in os.walk(temp_dir):
                for file in files:
                    if file != 'fixed.pptx':
                        file_path = os.path.join(root, file)
                        arcname = os.path.relpath(file_path, temp_dir)
                        zip_ref.write(file_path, arcname)
        
        # Create a copy in a persistent location
        persistent_path = os.path.join(tempfile.gettempdir(), f'fixed_{uuid.uuid4()}.pptx')
        shutil.copy2(new_pptx_path, persistent_path)
        
        return persistent_path
    except Exception as e:
        print(f"Error preprocessing PPTX file: {e}")
        return pptx_file
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)

def pptx_to_fabric_json(pptx_file):
    """Convert PPTX file to Fabric.js JSON format"""
    # Create a unique folder for this upload
    upload_id = datetime.now().strftime('%Y%m%d_%H%M%S_') + str(uuid.uuid4())[:8]
    image_folder = os.path.join('static', 'uploads', upload_id)
    os.makedirs(image_folder, exist_ok=True)

    processed_pptx = None
    try:
        # Preprocess the PPTX file
        processed_pptx = preprocess_pptx_file(pptx_file)
        
        # Load the presentation
        prs = Presentation(processed_pptx)
        
        # Initialize handlers
        color_handler = ColorHandler(prs)
        text_handler = TextHandler(color_handler)
        shape_handler = ShapeHandler(color_handler, text_handler)
        
        slides_data = []
        for slide_index, slide in enumerate(prs.slides):
            slide_objects = []
            
            # Process background if exists
            if hasattr(slide, 'background') and slide.background.fill:
                bg_color = color_handler.get_shape_color(slide.background)
                if bg_color:
                    slide_objects.append({
                        "type": "rect",
                        "left": 0,
                        "top": 0,
                        "width": prs.slide_width / 12700,
                        "height": prs.slide_height / 12700,
                        "fill": bg_color,
                        "selectable": False
                    })

            # Process all shapes
            for shape in slide.shapes:
                shape_data = shape_handler.process_shape(shape)
                if shape_data:
                    if shape_data["type"] == "group":
                        slide_objects.extend(shape_data["objects"])
                    else:
                        slide_objects.append(shape_data)

            slides_data.append({
                "objects": slide_objects,
                "width": prs.slide_width / 12700,
                "height": prs.slide_height / 12700,
                "slideNumber": slide_index + 1
            })

        return slides_data
        
    except Exception as e:
        print(f"Error processing PPTX file: {e}")
        raise
    finally:
        # Clean up the processed file
        if processed_pptx and processed_pptx != pptx_file:
            try:
                os.remove(processed_pptx)
            except Exception as e:
                print(f"Error cleaning up temporary file: {e}")
